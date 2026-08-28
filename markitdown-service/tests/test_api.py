from __future__ import annotations

import json
import os
import subprocess
import sys
from dataclasses import replace
from io import BytesIO
from pathlib import Path

from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation

from app import main as app_module

app_module.settings = replace(app_module.settings, api_key="test-api-key")
client = TestClient(app_module.app)
auth_headers = {"X-API-Key": "test-api-key"}


def test_health_is_public() -> None:
    response = client.get("/health")
    assert response.status_code == 200
    assert response.json()["status"] == "healthy"
    assert response.headers["cache-control"] == "no-store"


def test_conversion_requires_authentication() -> None:
    response = client.post(
        "/v1/convert",
        files={"file": ("example.txt", b"Hello", "text/plain")},
    )
    assert response.status_code == 401


def test_rejects_invalid_api_key() -> None:
    response = client.get("/v1/formats", headers={"X-API-Key": "wrong"})
    assert response.status_code == 401


def test_converts_plain_text_to_markdown() -> None:
    response = client.post(
        "/v1/convert",
        headers=auth_headers,
        files={
            "file": (
                "example.txt",
                "عنوان\n\nمحتوى عربي واضح".encode(),
                "text/plain",
            )
        },
    )
    assert response.status_code == 200, response.text
    body = response.json()
    assert body["filename"] == "example.txt"
    assert "محتوى عربي واضح" in body["markdown"]
    assert body["engine"] == "microsoft/markitdown@0.1.7"


def test_can_return_raw_markdown() -> None:
    response = client.post(
        "/v1/convert?response_format=markdown",
        headers=auth_headers,
        files={"file": ("example.md", b"# Heading", "text/markdown")},
    )
    assert response.status_code == 200
    assert response.headers["content-type"].startswith("text/markdown")
    assert "# Heading" in response.text


def test_converts_excel_workbook() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "العقود"
    sheet.append(["رقم العقد", "الحالة"])
    sheet.append([101, "ساري"])
    payload = BytesIO()
    workbook.save(payload)

    response = client.post(
        "/v1/convert",
        headers=auth_headers,
        files={
            "file": (
                "contracts.xlsx",
                payload.getvalue(),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    )
    assert response.status_code == 200, response.text
    markdown = response.json()["markdown"]
    assert "رقم العقد" in markdown
    assert "ساري" in markdown


def test_converts_powerpoint_presentation() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "ملخص المشروع"
    slide.placeholders[1].text = "خدمة مستقلة لجميع المشاريع"
    payload = BytesIO()
    presentation.save(payload)

    response = client.post(
        "/v1/convert",
        headers=auth_headers,
        files={
            "file": (
                "project.pptx",
                payload.getvalue(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert response.status_code == 200, response.text
    markdown = response.json()["markdown"]
    assert "ملخص المشروع" in markdown
    assert "خدمة مستقلة" in markdown


def test_rejects_unsupported_extensions() -> None:
    response = client.post(
        "/v1/convert",
        headers=auth_headers,
        files={"file": ("archive.zip", b"not-a-zip", "application/zip")},
    )
    assert response.status_code == 415


def test_shared_action_discovers_documents_without_converting_site_assets(
    tmp_path: Path,
) -> None:
    project = tmp_path / "sample-project"
    docs = project / "docs"
    docs.mkdir(parents=True)
    (docs / "notes.txt").write_text("نص مشروع سابق", encoding="utf-8")
    (project / "index.html").write_text("<h1>واجهة الموقع</h1>", encoding="utf-8")

    repository_root = Path(__file__).resolve().parents[2]
    converter = repository_root / ".github/actions/markitdown-batch/convert.py"
    environment = os.environ.copy()
    environment["GITHUB_WORKSPACE"] = str(tmp_path)
    completed = subprocess.run(
        [
            sys.executable,
            str(converter),
            "--source-dir",
            ".",
            "--output-dir",
            ".markitdown-output",
            "--mode",
            "documents",
            "--max-file-size-mb",
            "50",
            "--max-files",
            "10",
            "--max-output-characters",
            "5000000",
            "--timeout-seconds",
            "30",
            "--fail-on-error",
            "true",
        ],
        cwd=tmp_path,
        env=environment,
        capture_output=True,
        check=False,
        text=True,
        timeout=60,
    )

    assert completed.returncode == 0, completed.stderr
    output = tmp_path / ".markitdown-output"
    assert (output / "sample-project/docs/notes.txt.md").is_file()
    assert not (output / "sample-project/index.html.md").exists()
    manifest = json.loads((output / "manifest.json").read_text(encoding="utf-8"))
    assert len(manifest["converted"]) == 1
    assert manifest["errors"] == []
